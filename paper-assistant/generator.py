import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from config import OUTPUT_DIR, IMAGES_DIR

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT = RGBColor(0x1A, 0x56, 0xDB)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)


def create_presentation(paper_slides_list, output_name="presentation.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank_layout)
    _add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
                  "Group Meeting Report", Pt(44), ACCENT, bold=True, align=PP_ALIGN.LEFT)
    _add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.8),
                  f"Papers Covered: {len(paper_slides_list)}", Pt(22), GRAY, align=PP_ALIGN.LEFT)
    _add_divider(slide, Inches(1.5), Inches(4.0), Inches(10.3))

    for paper_data in paper_slides_list:
        slides = paper_data.get("slides", [])

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)

            title = slide_data.get("title", "")
            _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9),
                          title, Pt(30), ACCENT, bold=True, align=PP_ALIGN.LEFT)

            _add_divider(slide, Inches(0.8), Inches(1.3), Inches(11.7))

            bullets = slide_data.get("bullets", [])
            body_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(5.2))
            tf = body_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"  {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = DARK
                p.space_after = Pt(14)
                p.space_before = Pt(4)

        _embed_figures_slide(prs, paper_data, lambda: prs.slides.add_slide(blank_layout))

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, output_name)
    prs.save(output_path)
    return output_path


def append_to_presentation(paper_slides_list, existing_path, output_name):
    """Append new paper slides to an existing PPTX file."""
    prs = Presentation(existing_path)
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    # Separator slide to mark new content
    slide = prs.slides.add_slide(blank_layout)
    _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
                  f"+ {len(paper_slides_list)} New Paper(s)", Pt(40), ACCENT, bold=True, align=PP_ALIGN.CENTER)
    _add_divider(slide, Inches(3), Inches(4.0), Inches(7.3))

    for paper_data in paper_slides_list:
        slides = paper_data.get("slides", [])

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)

            title = slide_data.get("title", "")
            _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9),
                          title, Pt(30), ACCENT, bold=True, align=PP_ALIGN.LEFT)

            _add_divider(slide, Inches(0.8), Inches(1.3), Inches(11.7))

            bullets = slide_data.get("bullets", [])
            body_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(5.2))
            tf = body_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"  {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = DARK
                p.space_after = Pt(14)
                p.space_before = Pt(4)

        _embed_figures_slide(prs, paper_data, lambda: prs.slides.add_slide(blank_layout))

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, output_name)
    prs.save(output_path)
    return output_path


def _add_text_box(slide, left, top, width, height, text, font_size, color,
                  bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb


def _add_divider(slide, left, top, width):
    line = slide.shapes.add_shape(1, left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return line


def _embed_figures_slide(prs, paper_data, slide_func):
    """Add a figures slide embedding extracted images from the paper."""
    images = paper_data.get("images", [])
    if not images:
        return

    paper_id = str(paper_data.get("paper_id", ""))
    paper_title = paper_data.get("title", "")
    available = []
    for img in images:
        img_path = os.path.join(IMAGES_DIR, paper_id, img.get("filename", ""))
        if os.path.exists(img_path):
            available.append((img_path, img))

    if not available:
        return

    slide = slide_func()
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.7),
                  f"Figures — {paper_title[:60]}", Pt(24), ACCENT, bold=True)
    _add_divider(slide, Inches(0.8), Inches(1.0), Inches(11.7))

    # Layout images in a grid (max 4 per slide)
    positions = [
        (Inches(0.6), Inches(1.4), Inches(5.8), Inches(2.7)),
        (Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.7)),
        (Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.7)),
        (Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.7)),
    ]

    for i, (img_path, img) in enumerate(available[:4]):
        left, top, width, height = positions[i]
        try:
            pic = slide.shapes.add_picture(img_path, left, top, width, height)
        except Exception:
            continue

        # Add page label
        _add_text_box(slide, left, top + height + Inches(0.05), width, Inches(0.25),
                      f"Page {img.get('page', '?')}", Pt(9), GRAY, align=PP_ALIGN.CENTER)
